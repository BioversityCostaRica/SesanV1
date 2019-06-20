# coding=utf-8
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from get_vals import getDashReportData, dataReport

import os
from binascii import a2b_base64

from pyramid.response import FileResponse


# merge cells vertically
def mergeCellsVertically(table, start_row_idx, end_row_idx, col_idx):
    row_count = end_row_idx - start_row_idx + 1
    column_cells = [r.cells[col_idx] for r in table.rows][start_row_idx:]
    column_cells[0]._tc.set('rowSpan', str(row_count))


def genPPTX(self, date):
    # load a presentation
    path = os.path.join(self.request.registry.settings['user.repository'], "TMP")

    meses = ["Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio", "Julio", "Agosto", "Septiembre", "Octubre",
             "Noviembre", "Diciembre"]
    data = getDashReportData(self, str(meses.index(date[0]) + 1), date[1])

    data.pop('signatures', None)

    binary_data = a2b_base64(self.request.POST.get("pptx_img").replace("data:image/png;base64,", ""))
    fd = open(path + '/image1.png', 'wb')
    fd.write(binary_data)
    fd.close()

    binary_data = a2b_base64(self.request.POST.get("pptx_img2").replace("data:image/png;base64,", ""))
    fd = open(path + '/image2.png', 'wb')
    fd.write(binary_data)
    fd.close()

    binary_data = a2b_base64(self.request.POST.get("pptx_img3").replace("data:image/png;base64,", ""))
    fd = open(path + '/image3.png', 'wb')
    fd.write(binary_data)
    fd.close()

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]

    # slide 1
    slide = prs.slides.add_slide(title_slide_layout)

    txBox = slide.shapes.add_textbox(Cm(4.9), Cm(2), Cm(15), Cm(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = u"Situación de Seguridad Alimentaria y Nutricional\n COMUSAN %s, %s, %s" % (str(self.user.munic).decode("latin1"),date[0], date[1])
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(114, 159, 207)
    p.alignment = PP_ALIGN.CENTER

    table = slide.shapes.add_table(4, 2, Cm(2.3), Cm(10), Cm(20), Cm(2)).table
    table.first_row = False
    table.horz_banding = False
    table.vert_banding = True

    table.cell(0, 0).text = u"Fecha de elaboración del informe"
    table.cell(1, 0).text = u"Fuente de informacion:"
    table.cell(2, 0).text = u"Responsable del informe"
    table.cell(3, 0).text = u"Monitor"
    table.cell(0, 1).text = u"%s, %s" % (date[0], date[1])
    table.cell(1, 1).text = u"COMUSAN de %s" % str(self.user.munic).decode("latin1")
    table.cell(2, 1).text = u"Oficina Municipal de SAN / SESAN"
    table.cell(3, 1).text = u"%s" % str(self.user.fullName).decode("latin1")

    # slide 2
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    pic = slide.shapes.add_picture(self.request.registry.settings['user.repository']+"/TMP/image1.png", Cm(1), Cm(2), Cm(10))

    pic = slide.shapes.add_picture(self.request.registry.settings['user.repository']+"/TMP/image3.png", Cm(14), Cm(2), Cm(8))

    # slide 3
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    pic = slide.shapes.add_picture(self.request.registry.settings['user.repository']+"/TMP/image2.png", Cm(2.6), Cm(0.5), Cm(20))

    # slide 4


    rowcount = 0

    for p in data.keys():
        if "_alert" not in p and "date" not in p and "signatures" not in p and p not in ["comunidades", "acciones",
                                                                                         "coverage", "comunidades2",
                                                                                         "san"]:
            for d in data[p]:
                for v in data[p][d]["var"]:
                    rowcount += 1

    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    tab = slide.shapes.add_table(rowcount + 2, 8, Cm(1), Cm(1), Cm(24),
                                 Cm(10)).table  # rows, cols, left, top, width, heigh

    tab.first_row = False
    tab.horz_banding = False
    tab.vert_banding = True
    tab.cell(0, 0).text = u"Reporte Municipal de Situación General de SAN con base a la Sala Situacional "

    row_cells = [c for c in tab.rows[0].cells][0:2]
    row_cells[0]._tc.set('gridSpan', str(8))

    tab.cell(1, 0).text = "Indice de Situacion General de SAN"
    tab.cell(1, 1).text = "situacion actual"
    tab.cell(1, 2).text = "Indice de afectacion por pilar"
    tab.cell(1, 3).text = "Situacion actual "
    tab.cell(1, 4).text = "Indice de afectacion por indicador"
    tab.cell(1, 5).text = "Situacion actual"
    tab.cell(1, 6).text = "Afectacion por variable"
    tab.cell(1, 7).text = "Situacion actual"

    rowcount = 0

    for p in data.keys():
        if "_alert" not in p and "date" not in p and "signatures" not in p and p not in ["comunidades", "acciones",
                                                                                         "coverage", "comunidades2",
                                                                                         "san"]:

            init = rowcount + 2

            for d in data[p]:
                if len(data[p][d]["var"]) > 1:
                    #only work if len(data[p][d]["var"]) is 2 or 3, its necessary update this part for 4 or more
                    if len(data[p][d]["var"]) > 2:
                        idr = rowcount + len(data[p][d]["var"])-1
                    else:
                        idr = rowcount + len(data[p][d]["var"])
                else:
                    idr = rowcount + 2

                tab.cell(idr, 4).text = d.decode("latin1").replace("_", " ")  # x[p][d]["val"][0]
                tab.cell(idr, 5).text = str(float(data[p][d]["val"][0]))
                tab.cell(idr, 5).fill.solid()
                c = tuple(int(data[p][d]["val"][1][0][0].lstrip('#')[i:i + 2], 16) for i in (0, 2, 4))
                tab.cell(idr, 5).fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])


                mergeCellsVertically(tab, idr, idr + len(data[p][d]["var"]) - 1, 4)
                mergeCellsVertically(tab, idr, idr + len(data[p][d]["var"]) - 1, 5)
                for v in data[p][d]["var"]:
                    rowcount += 1

                    tab.cell(rowcount + 1, 6).text = v[0].decode("latin1").replace("_", " ")
                    tab.cell(rowcount + 1, 7).text = str(float(v[3]))# color
                    tab.cell(rowcount + 1, 7).fill.solid()
                    c= tuple(int(v[5][0].lstrip('#')[i:i + 2], 16) for i in (0, 2, 4))
                    tab.cell(rowcount + 1, 7).fill.fore_color.rgb = RGBColor(c[0],c[1],c[2])



            tab.cell(init, 2).text = p.decode("latin1").replace("_", " ")


            tab.cell(init, 3).text = str(float(data[p + "_alert"][0]))
            tab.cell(init, 3).fill.solid()
            c = tuple(int(data[p + "_alert"][2].lstrip('#')[i:i + 2], 16) for i in (0, 2, 4))
            tab.cell(init, 3).fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])

            mergeCellsVertically(tab, init, rowcount + 1, 2)
            mergeCellsVertically(tab, init, rowcount + 1, 3)

    san = dataReport(self, str(meses.index(date[0]) + 1), data["date"][1])
    print san
    data["san"] = san["san"]
    tab.cell(2, 0).text = data["san"][2]
    tab.cell(2, 1).text = str(float(data["san"][0]))
    tab.cell(2, 1).fill.solid()
    c = tuple(int(data["san"][1].lstrip('#')[i:i + 2], 16) for i in (0, 2, 4))
    tab.cell(2, 1).fill.fore_color.rgb = RGBColor(c[0], c[1], c[2])


    mergeCellsVertically(tab, 2, rowcount + 1, 0)
    mergeCellsVertically(tab, 2, rowcount + 1, 1)

    # aqui quede
    # amazing for merge cells in python-pptx
    # https://groups.google.com/forum/#!topic/python-pptx/cVRP9sSpEjA


    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

    for cell in iter_cells(tab):
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(5)

    prs.save(path + '/report.pptx')

    response = FileResponse(path + '/report.pptx', request=self.request)
    headers = response.headers
    headers['Content-Type'] = 'application/download'
    headers['Accept-Ranges'] = 'bite'
    headers['Content-Disposition'] = 'attachment;filename=' + "report.pptx"
    return response


    # genPPTX()

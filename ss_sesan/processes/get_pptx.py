# coding=utf-8
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor

import os
from binascii import a2b_base64

from pyramid.response import FileResponse


def genPPTX(self):
    # load a presentation
    path = os.path.join(self.request.registry.settings['user.repository'], "TMP")


    binary_data = a2b_base64(self.request.POST.get("pptx_img").replace("data:image/png;base64,",""))
    fd = open(path+'/image1.png', 'wb')
    fd.write(binary_data)
    fd.close()

    binary_data = a2b_base64(self.request.POST.get("pptx_img2").replace("data:image/png;base64,", ""))
    fd = open(path + '/image2.png', 'wb')
    fd.write(binary_data)
    fd.close()

    binary_data = a2b_base64(self.request.POST.get("pptx_img3").replace("data:image/png;base64,", ""))
    fd = open(path + '/image3.png', 'wb')
    fd.write(binary_data)
    fd.close()




    prs = Presentation()
    title_slide_layout = prs.slide_layouts[6]

    #slide 1
    slide = prs.slides.add_slide(title_slide_layout)

    txBox = slide.shapes.add_textbox(Cm(4.9),  Cm(2), Cm(15), Cm(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = u"Situación de Seguridad Alimentaria y Nutricional\n COMUSAN San José Del Golfo, Febrero,2018"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(114,159, 207)
    p.alignment = PP_ALIGN.CENTER

    table=slide.shapes.add_table(3, 2,Cm(2.3), Cm(10),Cm(20),Cm(2)).table
    table.first_row = False
    table.horz_banding = False
    table.vert_banding = True

    table.cell(0, 0).text = u"Fecha de elaboración del informe"
    table.cell(1, 0).text = u"Fuente de informacion:"
    table.cell(2, 0).text = u"Responsable del informe"
    table.cell(0, 1).text = u"Febrero,2018"
    table.cell(1, 1).text = u"COMUSAN de San José Del Golfo"
    table.cell(2, 1).text = u"Oficina Municipal de SAN / SESAN"


    #slide 2
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    pic =slide.shapes.add_picture("/home/acoto/sesan_repo/TMP/image1.png", Cm(1), Cm(2), Cm(10))

    pic = slide.shapes.add_picture("/home/acoto/sesan_repo/TMP/image3.png", Cm(14), Cm(2), Cm(8))

    #slide 3
    title_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_slide_layout)
    pic = slide.shapes.add_picture("/home/acoto/sesan_repo/TMP/image2.png", Cm(2.6), Cm(0.5), Cm(20))

    prs.save(path+'/report.pptx')

    response = FileResponse(path+'/report.pptx', request=self.request)
    headers = response.headers
    headers['Content-Type'] = 'application/download'
    headers['Accept-Ranges'] = 'bite'
    headers['Content-Disposition'] = 'attachment;filename=' + "report.pptx"
    return response


#genPPTX()